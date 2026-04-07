"""Shared document extraction helpers for uploads and PPT workflows."""

from __future__ import annotations

import logging
import re
import subprocess
from pathlib import Path

logger = logging.getLogger(__name__)


def _strip_yaml_frontmatter(text: str) -> str:
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            return text[end + 5 :].lstrip()
    return text


def _extract_markdown_tables(text: str) -> list[list[list[str]]]:
    tables: list[list[list[str]]] = []
    current: list[list[str]] = []
    for raw_line in text.split("\n"):
        line = raw_line.strip()
        if not line.startswith("|"):
            if current:
                tables.append(current)
                current = []
            continue
        if re.match(r"^\|[\s\-:|]+\|$", line):
            continue
        cells = [cell.strip() for cell in line.strip("|").split("|")]
        if cells:
            current.append(cells)
    if current:
        tables.append(current)
    return tables


def _table_to_markdown(table: list[list[str]]) -> list[str]:
    if not table:
        return []
    normalized_rows = [
        [str(cell or "").replace("\n", " ").strip() for cell in row]
        for row in table
        if any(str(cell or "").strip() for cell in row)
    ]
    if len(normalized_rows) < 2:
        return []
    header = normalized_rows[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * len(header)) + " |",
    ]
    for row in normalized_rows[1:]:
        padded = row + [""] * max(0, len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    return lines


def _run_markitdown(file_path: str) -> str:
    try:
        from markitdown import MarkItDown

        result = MarkItDown().convert(file_path)
        return (result.text_content or "").strip()
    except Exception as exc:
        logger.warning("[DocumentExtraction] markitdown SDK 提取失败: %s", exc)

    try:
        result = subprocess.run(
            ["python", "-m", "markitdown", file_path],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode == 0:
            return result.stdout.strip()
        logger.warning("[DocumentExtraction] markitdown CLI 失败: %s", result.stderr.strip())
    except Exception as exc:
        logger.warning("[DocumentExtraction] markitdown CLI 调用失败: %s", exc)

    return ""


def _ocr_pdf_to_text(file_path: str) -> str:
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError:
        return ""

    try:
        images = convert_from_path(file_path, dpi=200)
    except Exception as exc:
        logger.warning("[DocumentExtraction] OCR 渲染 PDF 失败: %s", exc)
        return ""

    text_parts: list[str] = []
    for image in images:
        try:
            text = pytesseract.image_to_string(image)
        except Exception as exc:
            logger.warning("[DocumentExtraction] OCR 识别单页失败: %s", exc)
            continue
        if text and text.strip():
            text_parts.append(text.strip())
    return "\n\n".join(text_parts).strip()


def extract_text_from_pdf(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    text_parts: list[str] = []
    tables: list[list[list[str]]] = []
    page_count = 0

    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                text = (page.extract_text() or "").strip()
                if text:
                    text_parts.append(text)

                page_tables = page.extract_tables() or []
                if not page_tables:
                    page_tables = (
                        page.extract_tables(
                            table_settings={
                                "vertical_strategy": "lines",
                                "horizontal_strategy": "lines",
                            }
                        )
                        or []
                    )
                for table in page_tables:
                    if table and len(table) >= 2:
                        tables.append(table)
        if text_parts or tables:
            return "\n\n".join(text_parts).strip(), tables, page_count
    except ImportError:
        logger.warning("[DocumentExtraction] pdfplumber 不可用，回退到其他方案")
    except Exception as exc:
        logger.warning("[DocumentExtraction] pdfplumber 提取失败: %s", exc)

    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        page_count = max(page_count, len(reader.pages))
        for page in reader.pages:
            text = (page.extract_text() or "").strip()
            if text:
                text_parts.append(text)
        if text_parts:
            return "\n\n".join(text_parts).strip(), tables, page_count
    except Exception as exc:
        logger.warning("[DocumentExtraction] pypdf 提取失败: %s", exc)

    ocr_text = _ocr_pdf_to_text(file_path)
    if ocr_text:
        return ocr_text, tables, page_count

    markdown_text = _run_markitdown(file_path)
    if markdown_text:
        return markdown_text, _extract_markdown_tables(markdown_text), page_count

    return "", tables, page_count


def extract_text_from_docx(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        tables: list[list[list[str]]] = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                tables.append(rows)
        if paragraphs or tables:
            return "\n\n".join(paragraphs).strip(), tables, max(1, len(paragraphs) // 30)
    except ImportError:
        logger.warning("[DocumentExtraction] python-docx 不可用，尝试其他方案")
    except Exception as exc:
        logger.warning("[DocumentExtraction] python-docx 提取失败: %s", exc)

    try:
        result = subprocess.run(
            ["pandoc", file_path, "-o", "-", "--to", "plain"],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip(), [], 0
    except Exception as exc:
        logger.warning("[DocumentExtraction] pandoc 提取失败: %s", exc)

    markdown_text = _run_markitdown(file_path)
    if markdown_text:
        return markdown_text, _extract_markdown_tables(markdown_text), 0

    return "", [], 0


def extract_text_from_markdown(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    text = ""
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            text = Path(file_path).read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if not text:
        logger.error("[DocumentExtraction] Markdown 读取失败")
        return "", [], 0

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _strip_yaml_frontmatter(text)
    return text.strip(), _extract_markdown_tables(text), max(1, len(text) // 3000)


def extract_text_from_pptx(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = paragraph.text.strip()
                        if line:
                            text_parts.append(line)
                if shape.has_table:
                    rows: list[list[str]] = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        tables.append(rows)
        if text_parts or tables:
            return "\n\n".join(text_parts).strip(), tables, len(prs.slides)
    except ImportError:
        logger.warning("[DocumentExtraction] python-pptx 不可用，尝试 markitdown")
    except Exception as exc:
        logger.warning("[DocumentExtraction] python-pptx 提取失败: %s", exc)

    markdown_text = _run_markitdown(file_path)
    if markdown_text:
        return markdown_text, _extract_markdown_tables(markdown_text), 0

    return "", [], 0


def extract_document_content(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    if suffix in {".doc", ".docx"}:
        return extract_text_from_docx(file_path)
    if suffix == ".md":
        return extract_text_from_markdown(file_path)
    if suffix in {".ppt", ".pptx"}:
        return extract_text_from_pptx(file_path)
    raise ValueError(f"Unsupported file format: {suffix}")


def format_extracted_markdown(
    raw_text: str,
    tables: list[list[list[str]]] | None = None,
    *,
    source_name: str | None = None,
    page_count: int | None = None,
) -> str:
    lines: list[str] = []
    if source_name:
        lines.append(f"# {source_name}")
        lines.append("")
    if page_count:
        lines.append(f"- 页数: {page_count}")
        lines.append("")

    cleaned_text = (raw_text or "").strip()
    if cleaned_text:
        lines.append("## 提取文本")
        lines.append("")
        lines.append(cleaned_text)
        lines.append("")

    for index, table in enumerate(tables or [], start=1):
        table_lines = _table_to_markdown(table)
        if not table_lines:
            continue
        lines.append(f"## 表格 {index}")
        lines.append("")
        lines.extend(table_lines)
        lines.append("")

    return "\n".join(lines).strip() + "\n"
