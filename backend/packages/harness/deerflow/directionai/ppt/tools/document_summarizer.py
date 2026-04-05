"""Document summarizer utilities for upload-driven PPT generation."""

from __future__ import annotations

import json
import logging
import re
import subprocess
from pathlib import Path
from typing import Any, Callable

from openai import OpenAI
from pydantic import BaseModel, Field, field_validator

from .. import config as ppt_config
from .openai_compat import build_chat_completion_kwargs, stream_chat_completion_text

logger = logging.getLogger(__name__)

_CONTAINER_SKILLS_ROOT = Path("/mnt/skills/public")
_REPO_SKILLS_ROOT = Path(__file__).resolve().parents[7] / "skills" / "public"


def _resolve_skill_dir(name: str) -> Path:
    for root in (_CONTAINER_SKILLS_ROOT, _REPO_SKILLS_ROOT):
        candidate = root / name
        if candidate.exists():
            return candidate
    return _CONTAINER_SKILLS_ROOT / name


SKILL_ROOT = _resolve_skill_dir("document-summarizer")
PROCESSOR_ROOTS = {
    "pdf_processor": _resolve_skill_dir("document-processor-pdf"),
    "docx_processor": _resolve_skill_dir("document-processor-docx"),
    "markdown_processor": _resolve_skill_dir("document-processor-markdown"),
    "pptx_processor": _resolve_skill_dir("document-processor-pptx"),
}


def skill_paths() -> dict[str, str]:
    return {
        "document_summarizer": str(SKILL_ROOT / "SKILL.md"),
        **{name: str(path / "SKILL.md") for name, path in PROCESSOR_ROOTS.items()},
    }


def assert_skills_present() -> None:
    for name, path in skill_paths().items():
        if not Path(path).exists():
            raise FileNotFoundError(f"Skill 文件缺失: {name} → {path}")


class TableSpec(BaseModel):
    description: str
    headers: list[str]
    rows: list[list[str]]


class SectionSummary(BaseModel):
    section_title: str
    section_summary: str
    key_points: list[str] = Field(default_factory=list)
    tables: list[TableSpec] = Field(default_factory=list)
    important_data: list[str] = Field(default_factory=list)
    suggested_slide_count: int = Field(ge=1, le=10)
    narrative_transition: str | None = None

    @field_validator("key_points", mode="before")
    @classmethod
    def _ensure_list(cls, value: Any) -> list[str]:
        if isinstance(value, list):
            return [str(item) for item in value]
        return []


class Metadata(BaseModel):
    source_file: str | None = None
    total_pages: int | None = None
    language: str = "auto-detect"
    estimated_reading_time_minutes: int | None = None


class PPTGenerationHints(BaseModel):
    suggested_total_slides: int = Field(ge=4, le=20, default=8)
    audience: str = "general"
    style_preference: str = "informative"
    recommended_visual_elements: list[str] = Field(default_factory=list)
    content_focus: str = "knowledge_sharing"
    special_requirements: str | None = None


class DocumentSummary(BaseModel):
    document_title: str
    document_summary: str
    sections: list[SectionSummary] = Field(min_length=1)
    metadata: Metadata = Field(default_factory=Metadata)
    ppt_generation_hints: PPTGenerationHints = Field(default_factory=PPTGenerationHints)

    def to_planner_context(self) -> str:
        lines = [
            f"# {self.document_title}",
            f"## 概述\n{self.document_summary}",
            "",
        ]

        for index, section in enumerate(self.sections, start=1):
            lines.append(f"## {index}. {section.section_title}")
            lines.append(f"**本节概要**: {section.section_summary}")
            if section.key_points:
                lines.append("**要点**:")
                for point in section.key_points:
                    lines.append(f"  - {point}")
            if section.important_data:
                lines.append("**关键数据**:")
                for data in section.important_data:
                    lines.append(f"  - {data}")
            for table in section.tables:
                lines.append(f"**表格**: {table.description}")
                if table.headers:
                    lines.append("  | " + " | ".join(table.headers) + " |")
                    lines.append("  | " + " | ".join(["---"] * len(table.headers)) + " |")
                for row in table.rows[:5]:
                    lines.append("  | " + " | ".join(str(cell) for cell in row) + " |")
                if len(table.rows) > 5:
                    lines.append("  | ... |")
            if section.narrative_transition:
                lines.append(f"*→ 下一节: {section.narrative_transition}*")
            lines.append("")

        hints = self.ppt_generation_hints
        lines.append("## PPT 生成建议")
        lines.append(f"- 建议页数: {hints.suggested_total_slides}")
        lines.append(f"- 目标受众: {hints.audience}")
        lines.append(f"- 风格偏好: {hints.style_preference}")
        if hints.recommended_visual_elements:
            lines.append(f"- 推荐视觉元素: {', '.join(hints.recommended_visual_elements)}")
        if hints.content_focus:
            lines.append(f"- 内容类型: {hints.content_focus}")
        if hints.special_requirements:
            lines.append(f"- 特殊要求: {hints.special_requirements}")
        return "\n".join(lines)


def extract_text_from_pdf(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    try:
        import pdfplumber

        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text.strip())
                for table in page.extract_tables():
                    if table and len(table) >= 2:
                        tables.append(table)
            return "\n\n".join(text_parts), tables, len(pdf.pages)
    except ImportError:
        logger.warning("[DocumentSummarizer] pdfplumber 不可用，回退到 pypdf")
    except Exception as exc:
        logger.warning("[DocumentSummarizer] pdfplumber 提取失败: %s", exc)

    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text.strip())
        return "\n\n".join(texts), [], len(reader.pages)
    except Exception as exc:
        logger.error("[DocumentSummarizer] pypdf 提取失败: %s", exc)
        return "", [], 0


def extract_text_from_docx(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        tables: list[list[list[str]]] = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                tables.append(rows)
        return "\n\n".join(paragraphs), tables, max(1, len(paragraphs) // 30)
    except ImportError:
        logger.warning("[DocumentSummarizer] python-docx 不可用，尝试 pandoc")
    except Exception as exc:
        logger.warning("[DocumentSummarizer] python-docx 提取失败: %s", exc)

    try:
        result = subprocess.run(
            ["pandoc", file_path, "-o", "-", "--to", "plain"],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0:
            return result.stdout, [], 0
    except Exception as exc:
        logger.warning("[DocumentSummarizer] pandoc 提取失败: %s", exc)

    return "", [], 0


def _strip_yaml_frontmatter(text: str) -> str:
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            return text[end + 5 :].lstrip()
    return text


def _extract_markdown_tables(text: str) -> list[list[list[str]]]:
    tables: list[list[list[str]]] = []
    current: list[list[str]] = []
    for raw_line in text.split("\n"):
        line = raw_line.strip()
        if not line.startswith("|"):
            if current:
                tables.append(current)
                current = []
            continue
        if re.match(r"^\|[\s\-:|]+\|$", line):
            continue
        cells = [cell.strip() for cell in line.strip("|").split("|")]
        if cells:
            current.append(cells)
    if current:
        tables.append(current)
    return tables


def extract_text_from_markdown(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    text = ""
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            text = Path(file_path).read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if not text:
        logger.error("[DocumentSummarizer] Markdown 读取失败")
        return "", [], 0

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _strip_yaml_frontmatter(text)
    return text.strip(), _extract_markdown_tables(text), max(1, len(text) // 3000)


def extract_text_from_pptx(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = paragraph.text.strip()
                        if line:
                            text_parts.append(line)
                if shape.has_table:
                    rows: list[list[str]] = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        tables.append(rows)
        return "\n\n".join(text_parts), tables, max(1, len(prs.slides))
    except ImportError:
        logger.warning("[DocumentSummarizer] python-pptx 不可用，尝试 markitdown")
    except Exception as exc:
        logger.warning("[DocumentSummarizer] python-pptx 提取失败: %s", exc)

    try:
        result = subprocess.run(
            ["python", "-m", "markitdown", file_path],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and result.stdout:
            return result.stdout.strip(), [], 0
    except Exception as exc:
        logger.warning("[DocumentSummarizer] markitdown 提取失败: %s", exc)
    return "", [], 0


def extract_document_content(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    if suffix in {".doc", ".docx"}:
        return extract_text_from_docx(file_path)
    if suffix == ".md":
        return extract_text_from_markdown(file_path)
    if suffix in {".ppt", ".pptx"}:
        return extract_text_from_pptx(file_path)
    raise ValueError(f"Unsupported file format: {suffix}")


def _load_skill() -> str:
    skill_md = SKILL_ROOT / "SKILL.md"
    if not skill_md.exists():
        raise FileNotFoundError(f"SKILL.md not found at {skill_md}")
    return skill_md.read_text(encoding="utf-8")


def _build_messages(
    raw_text: str,
    tables: list[list[list[str]]],
    metadata: dict[str, Any],
    skill_md: str,
) -> list[dict[str, str]]:
    user_parts = [
        "以下是从文档中提取的原始内容，请按 SKILL.md 规范输出结构化摘要 JSON。",
        "",
    ]
    if metadata.get("source_file"):
        user_parts.append(f"来源文件: {metadata['source_file']}")
    if metadata.get("page_count"):
        user_parts.append(f"文档页数: {metadata['page_count']}")
    user_parts.extend(["", "## 原始文本内容", raw_text])

    if tables:
        user_parts.extend(["", "## 提取的表格"])
        for index, table in enumerate(tables, start=1):
            if not table or len(table) < 2:
                continue
            user_parts.append(f"### 表格 {index}")
            for row in table[:10]:
                user_parts.append(" | ".join(str(cell) for cell in row))
            if len(table) > 10:
                user_parts.append("...（后续行已截断）")
            user_parts.append("")

    return [
        {"role": "system", "content": skill_md},
        {"role": "user", "content": "\n".join(user_parts)},
    ]


def _extract_json_blob(text: str) -> str:
    if not text or text[0] not in "{[":
        return text
    opener = text[0]
    closer = "}" if opener == "{" else "]"
    stack = [closer]
    in_string = False
    escape = False
    for index in range(1, len(text)):
        char = text[index]
        if escape:
            escape = False
            continue
        if char == "\\":
            escape = True
            continue
        if char == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if char == closer and stack and char == stack[-1]:
            stack.pop()
            if not stack:
                return text[: index + 1]
        elif char == opener:
            stack.append(closer)
    return text


def _repair_json(text: str) -> str:
    text = re.sub(r",(\s*[}\]])", r"\1", text)
    result = list(text)
    stack: list[str] = []
    in_string = False
    escape = False
    for char in text:
        if escape:
            escape = False
            continue
        if char == "\\":
            escape = True
            continue
        if char == '"':
            in_string = not in_string
            continue
        if not in_string:
            if char in "{[":
                stack.append("}" if char == "{" else "]")
            elif char in "}]" and stack and stack[-1] == char:
                stack.pop()
    while stack:
        result.append(stack.pop())
    return "".join(result)


def _parse_summary(raw_content: str) -> DocumentSummary:
    cleaned = str(raw_content or "").strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    positions = [(cleaned.find("{"), "{"), (cleaned.find("["), "[")]
    positions = [(pos, char) for pos, char in positions if pos >= 0]
    if positions:
        start, _ = min(positions, key=lambda item: item[0])
        blob = _extract_json_blob(cleaned[start:])
        if blob:
            cleaned = blob

    try:
        data = json.loads(cleaned, strict=False)
    except json.JSONDecodeError:
        repaired = _repair_json(cleaned)
        try:
            data = json.loads(repaired, strict=False)
        except json.JSONDecodeError as exc:
            raise ValueError(f"无法解析 LLM 返回的 JSON，内容前300字：{cleaned[:300]}") from exc

    data.setdefault("metadata", {})
    data.setdefault("ppt_generation_hints", {})
    return DocumentSummary.model_validate(data)


def summarize_document(
    raw_text: str,
    *,
    tables: list[list[list[str]]] | None = None,
    source_file: str | None = None,
    page_count: int | None = None,
    model_provider: str = "minmax",
    thinking_callback: Callable[[str], None] | None = None,
    max_input_chars: int = 80000,
) -> DocumentSummary:
    if not raw_text or not raw_text.strip():
        raise ValueError("原始文本为空，请先提取文档内容")

    truncated = False
    if len(raw_text) > max_input_chars:
        raw_text = raw_text[:max_input_chars]
        truncated = True

    skill_md = _load_skill()
    provider_settings = ppt_config.get_llm_provider_settings(model_provider)
    client = OpenAI(
        api_key=provider_settings["api_key"],
        base_url=provider_settings["base_url"],
    )
    messages = _build_messages(
        raw_text,
        tables or [],
        {
            "source_file": source_file,
            "page_count": page_count,
        },
        skill_md,
    )

    if truncated:
        messages[1] = {
            "role": "user",
            "content": messages[1]["content"]
            + f"\n\n[注意：原始文本已截断至前 {max_input_chars} 字符，请基于已有内容生成摘要，并在 special_requirements 中注明。]",
        }

    raw_content, reasoning = stream_chat_completion_text(
        client,
        model=provider_settings["model"],
        max_tokens=16384,
        messages=messages,
        on_reasoning_chunk=thinking_callback,
        **build_chat_completion_kwargs(provider_settings["model"]),
    )
    logger.debug("[DocumentSummarizer] reasoning 前100字：%s", reasoning[:100])
    summary = _parse_summary(raw_content)
    logger.info(
        "[DocumentSummarizer] 摘要生成完成: title=%r, sections=%s, suggested_slides=%s",
        summary.document_title,
        len(summary.sections),
        summary.ppt_generation_hints.suggested_total_slides,
    )
    return summary
